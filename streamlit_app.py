from pptx import Presentation

class PPTReviewer:
    def __init__(self, review_criteria):
        self.review_criteria = review_criteria  # 审查标准字典，例如：{'Content Structure': ...}

    def load_presentation(self, path):
        """加载 PowerPoint 文件"""
        try:
            self.prs = Presentation(path)
            return True
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return False

    def review_slide(self, slide_number=None):
        """
        审查指定幻灯片或所有幻灯片。
        Args:
            slide_number (int): 要审查的幻灯片编号（从0开始）。
        Returns:
            dict: 每个标准的评分结果和反馈。
        """
        if not self.prs:
            return {"error": "Presentation未加载"}

        slides = self.prs.slides
        total_slides = len(slides)

        # 如果 slide_number 为 None，则审查所有幻灯片；否则只审查指定幻灯片
        target_slides = [slides[i] for i in range(total_slides)] if slide_number is None else [slides[slide_number]]

        results = []
        for slide in target_slides:
            feedback = {}
            # 按照每个标准进行评分和反馈
            for criterion, weight in self.review_criteria.items():
                score = 0
                comment = "未评估"
                
                if criterion == 'Content Structure':
                    # 示例：检查标题是否存在
                    if slide.shapes.title is not None:
                        score += 5
                        comment = "幻灯片有清晰的标题。"
                    else:
                        comment = "建议添加一个明确的标题。"

                elif criterion == 'Staff Work':
                    # 示例：查找作者信息（假设在备注中）
                    notes = slide.notes.text if slide.notes is not None else ""
                    if "作者：" in notes or "Contributor:" in notes:
                        score += 4
                        comment = "幻灯片有明确的作者信息。"
                    else:
                        comment = "建议添加作者信息。"

                elif criterion == 'Clarity and Choice of Words':
                    # 示例：检查文本框中的内容是否简洁（字数限制）
                    text_boxes = [shape.text for shape in slide.shapes if hasattr(shape, 'text')]
                    total_words = sum(len(tb.split()) for tb in text_boxes)
                    if total_words <= 50:
                        score += 4
                        comment = "幻灯片文本简洁明了。"
                    else:
                        comment = f"建议减少文字，当前有 {total_words} 字。"

                # 添加更多审查标准的逻辑...

                feedback[criterion] = {
                    'score': int(score),
                    'comment': comment,
                    'weight': weight
                }

            results.append(feedback)

        return self._calculate_total_score(results)

    def _calculate_total_score(self, review_results):
        """
        计算总分。
        Args:
            review_results (list): 每个标准的评分结果和反馈列表。
        Returns:
            dict: 总结报告，包括每个幻灯片的标准得分和总体分数。
        """
        total_scores = []
        for slide_reviews in review_results:
            slide_score = 0
            for criterion, data in slide_reviews.items():
                slide_score += (data['score'] * data['weight'])
            total_scores.append(slide_score)

        average_score = sum(total_scores) / len(total_scores)
        return {
            'slide_scores': total_scores,
            'average_score': round(average_score, 2),
            'total_slides_reviewed': len(review_results)
        }

# 示例用法
if __name__ == "__main__":
    # 定义审查标准（权重可以根据需要调整）
    review_criteria = {
        'Content Structure': 0.15,
        'Staff Work': 0.1,
        'Clarity and Choice of Words': 0.2,
        'Strategy': 0.15,
        'Outcome': 0.1,
        'Descriptors and Analyses': 0.15,
        'Recommendations': 0.05
    }

    # 初始化工具
    reviewer = PPTReviewer(review_criteria)

    # 加载PPT文件
    if reviewer.load_presentation("example.pptx"):
        # 审查所有幻灯片（slide_number=None）或指定幻灯片（例如 slide_number=2）
        review_results = reviewer.review_slide(slide_number=None)
        
        print("\n审查结果：")
        for i, result in enumerate(review_results):
            print(f"\n幻灯片 {i + 1}: 平均分={result['average_score']}")
            # 打印每个标准的评分和反馈
            for criterion, data in result.items():
                if criterion != 'average_score':
                    print(f"{criterion}：得分={data['score']}，权重={data['weight']}，评论：{data['comment']}")
        print("\n总体平均分:", review_results[-1]['average_score'])
    else:
        exit("无法加载PPT文件。")
